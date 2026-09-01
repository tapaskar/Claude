"""Programmatic text-fit QA: measure every text box against its frame.

LibreOffice cannot load any pptx in this sandbox, so instead of eyeballing a
render we measure. The deck is set in IBM Plex Sans / IBM Plex Mono and the real
faces are installed here, so these measurements are exact rather than approximate
- provided the recipient also has IBM Plex Sans. If they do not, PowerPoint
substitutes and the fit is no longer guaranteed.
"""
from pptx import Presentation
from PIL import ImageFont

LIB = "/usr/share/fonts/truetype/liberation/"
PLEX = "/usr/share/fonts/truetype/plex/"
FACE = {
    ("IBM Plex Sans", False): PLEX + "IBMPlexSans-Regular.ttf",
    ("IBM Plex Sans", True):  PLEX + "IBMPlexSans-SemiBold.ttf",
    ("IBM Plex Mono", False): PLEX + "IBMPlexMono-Regular.ttf",
    ("IBM Plex Mono", True):  PLEX + "IBMPlexMono-Bold.ttf",
    ("Arial", False): LIB + "LiberationSans-Regular.ttf",
    ("Arial", True):  LIB + "LiberationSans-Bold.ttf",
    ("Calibri", False): LIB + "LiberationSans-Regular.ttf",
    ("Calibri", True):  LIB + "LiberationSans-Bold.ttf",
    ("Courier New", False): LIB + "LiberationMono-Regular.ttf",
    ("Courier New", True):  LIB + "LiberationMono-Bold.ttf",
}
DEFAULT_FACE = PLEX + "IBMPlexSans-Regular.ttf"

PX, EMU_IN = 4, 914400
_cache = {}


def font(name, sz, bold):
    key = (name, sz, bold)
    if key not in _cache:
        _cache[key] = ImageFont.truetype(
            FACE.get((name, bold), DEFAULT_FACE), int(round(sz * PX)))
    return _cache[key]


def w_pt(txt, spec):
    return font(*spec).getlength(txt) / PX


def para_segments(para):
    """[(word, (face,size,bold)), ...] preserving each run's own font."""
    segs = []
    for r in para.runs:
        if not r.text:
            continue
        spec = (r.font.name or "Calibri",
                r.font.size.pt if r.font.size else 18,
                bool(r.font.bold))
        # keep spaces attached so wrapping is word-accurate across run boundaries
        for i, word in enumerate(r.text.split(" ")):
            if word:
                segs.append((word, spec))
    return segs


def wrap(segs, box_pt):
    """Greedy wrap over mixed-font segments. Returns (lines, widest_single_word)."""
    if not segs:
        return 1, 0.0, None
    lines, cur, widest, worst = 1, 0.0, 0.0, None
    space = None
    for word, spec in segs:
        ww = w_pt(word, spec)
        if ww > widest:
            widest, worst = ww, word
        sp = w_pt(" ", spec) if cur > 0 else 0.0
        if cur + sp + ww <= box_pt or cur == 0:
            cur += sp + ww
        else:
            lines += 1
            cur = ww
    return lines, widest, worst


def line_h(para, segs):
    ls = para.line_spacing
    if ls is not None:
        # pptxgenjs writes a point value; python-pptx returns Length for spcPts
        return ls.pt if hasattr(ls, "pt") else ls * max(s[1][1] for s in segs) * 1.2
    return max(s[1][1] for s in segs) * 1.22 if segs else 12


issues = []
pr = Presentation("Vi-AI-NOC-Sizing.pptx")
for si, slide in enumerate(pr.slides, 1):
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        bw = (sh.width / EMU_IN) * 72
        bh = (sh.height / EMU_IN) * 72
        total, widest, worst = 0.0, 0.0, None
        for para in sh.text_frame.paragraphs:
            segs = para_segments(para)
            if not segs:
                total += 12
                continue
            bullet_indent = 16 if para._p.find(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr') is not None \
                and 'buChar' in para._p.xml else 0
            n, wd, wt = wrap(segs, bw - bullet_indent)
            if wd > widest:
                widest, worst = wd, wt
            total += n * line_h(para, segs) + (para.space_after.pt if para.space_after else 0)
        txt = sh.text_frame.text[:56].replace("\n", " / ")
        if total > bh + 1.0:
            issues.append((si, "OVERFLOW", f"needs {total:.0f}pt in {bh:.0f}pt", txt))
        if widest > bw + 1.0:
            issues.append((si, "TOO-NARROW", f"'{worst}' {widest:.0f}pt > {bw:.0f}pt", txt))
        if sh.left + sh.width > pr.slide_width + 1000:
            issues.append((si, "OFF-SLIDE", "past right edge", txt))
        if sh.top + sh.height > pr.slide_height + 1000:
            issues.append((si, "OFF-SLIDE", "past bottom edge", txt))

print(f"{len(issues)} issue(s)" if issues else "No text-fit issues found.")
for si, kind, det, txt in issues:
    print(f"  s{si} [{kind}] {det}\n        \"{txt}\"")
