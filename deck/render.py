"""Render the pptx shape tree to PNG with PIL.

Not a PowerPoint-accurate renderer - it draws each shape and text run at the exact
geometry authored in the file. That is enough to catch the defects visual QA is for:
overlap, collisions, uneven spacing, contrast, balance.
"""
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

LIB = "/usr/share/fonts/truetype/liberation/"
PLEX = "/usr/share/fonts/truetype/plex/"
FACE = {
    ("IBM Plex Sans", False): PLEX + "IBMPlexSans-Regular.ttf",
    ("IBM Plex Sans", True):  PLEX + "IBMPlexSans-SemiBold.ttf",
    ("IBM Plex Mono", False): PLEX + "IBMPlexMono-Regular.ttf",
    ("IBM Plex Mono", True):  PLEX + "IBMPlexMono-Bold.ttf",
    ("Arial", False): LIB + "LiberationSans-Regular.ttf",
    ("Arial", True):  LIB + "LiberationSans-Bold.ttf",
    ("Calibri", False): LIB + "LiberationSans-Regular.ttf",
    ("Calibri", True):  LIB + "LiberationSans-Bold.ttf",
    ("Courier New", False): LIB + "LiberationMono-Regular.ttf",
    ("Courier New", True):  LIB + "LiberationMono-Bold.ttf",
}
DEFAULT_FACE = PLEX + "IBMPlexSans-Regular.ttf"

E, DPI = 914400, 110
_fc = {}
def fnt(n,s,b):
    k=(n,round(s,1),b)
    if k not in _fc:
        _fc[k]=ImageFont.truetype(FACE.get((n,bool(b)),DEFAULT_FACE),
                                  max(6,int(round(s*DPI/72))))
    return _fc[k]
px = lambda emu: int(round(emu/E*DPI))

def solid(sh):
    try:
        f = sh.fill
        if f.type is not None and f.type == 1:
            return "#" + str(f.fore_color.rgb)
    except Exception:
        pass
    return None

def linecol(sh):
    try:
        c = sh.line.color
        if c and c.type is not None:
            return "#" + str(c.rgb)
    except Exception:
        pass
    return None

pr = Presentation("Vi-AI-NOC-Sizing.pptx")
W, H = px(pr.slide_width), px(pr.slide_height)

for si, slide in enumerate(pr.slides, 1):
    bg = "#FFFFFF"
    try:
        if slide.background.fill.type == 1:
            bg = "#" + str(slide.background.fill.fore_color.rgb)
    except Exception:
        pass
    img = Image.new("RGB", (W, H), bg)
    d = ImageDraw.Draw(img)

    for sh in slide.shapes:
        x, y, w, h = px(sh.left), px(sh.top), px(sh.width), px(sh.height)
        fill, ln = solid(sh), linecol(sh)
        st = str(getattr(sh, "shape_type", "")) or ""
        if fill or ln:
            if "ROUNDED" in st.upper():
                d.rounded_rectangle([x, y, x+w, y+h], radius=max(2, int(0.05*DPI)),
                                    fill=fill, outline=ln, width=1)
            elif "ARROW" in st.upper():
                d.polygon([(x,y+h*0.3),(x+w*0.6,y+h*0.3),(x+w*0.6,y),(x+w,y+h*0.5),
                           (x+w*0.6,y+h),(x+w*0.6,y+h*0.7),(x,y+h*0.7)], fill=fill or "#888")
            else:
                d.rectangle([x, y, x+w, y+h], fill=fill, outline=ln, width=1)

        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        cy = y
        for para in sh.text_frame.paragraphs:
            runs = [r for r in para.runs if r.text]
            if not runs:
                cy += 10; continue
            segs = []
            for r in runs:
                nm = r.font.name or "Calibri"
                sz = r.font.size.pt if r.font.size else 18
                bo = bool(r.font.bold)
                try: col = "#" + str(r.font.color.rgb)
                except Exception: col = "#000000"
                for wd in r.text.split(" "):
                    if wd: segs.append((wd, nm, sz, bo, col))
            if not segs: continue
            ls = para.line_spacing
            lh = (ls.pt if hasattr(ls,"pt") else None) or max(s[2] for s in segs)*1.22
            lh = int(round(lh*DPI/72))
            bullet = 'buChar' in para._p.xml
            indent = int(0.16*DPI) if bullet else 0
            align = str(para.alignment or "")
            # greedy wrap
            lines, cur = [], []
            cw = 0
            for seg in segs:
                sw = fnt(seg[1],seg[2],seg[3]).getlength(seg[0])
                sp = fnt(seg[1],seg[2],seg[3]).getlength(" ") if cur else 0
                if cur and cw+sp+sw > w-indent:
                    lines.append(cur); cur=[seg]; cw=sw
                else:
                    cur.append(seg); cw += sp+sw
            if cur: lines.append(cur)
            for li, line in enumerate(lines):
                tw = sum(fnt(s[1],s[2],s[3]).getlength(s[0]) for s in line) + \
                     sum(fnt(s[1],s[2],s[3]).getlength(" ") for s in line[1:])
                cx = x + indent
                if "RIGHT" in align.upper(): cx = x + w - tw
                elif "CENTER" in align.upper(): cx = x + (w-tw)/2
                if li == 0 and bullet:
                    d.ellipse([x+4, cy+lh*0.42, x+8, cy+lh*0.42+4], fill=line[0][4])
                base = cy + lh - max(s[2] for s in line)*DPI/72*1.02
                for s in line:
                    f = fnt(s[1],s[2],s[3])
                    d.text((cx, base), s[0], font=f, fill=s[4])
                    cx += f.getlength(s[0]) + f.getlength(" ")
                cy += lh
            if para.space_after: cy += int(para.space_after.pt*DPI/72)
    img.save(f"render-{si}.png")
    print("render-%d.png %dx%d" % (si, W, H))
